import os
from typing import Tuple

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB

SUPPORTED = {
    ".docx": "document",
    ".pdf": "document",
    ".txt": "document",
    ".pptx": "presentation",
    ".mp4": "video",
    ".mov": "video",
    ".m4v": "video",
    ".avi": "video",
    ".mp3": "audio",
    ".m4a": "audio",
    ".wav": "audio",
    ".ogg": "audio",
}


class FileProcessor:
    def detect_type(self, filename: str) -> str:
        ext = os.path.splitext(filename.lower())[1]
        return SUPPORTED.get(ext, "unknown")

    def validate_file(self, filename: str, size_bytes: int) -> Tuple[bool, str]:
        if size_bytes > MAX_FILE_SIZE:
            mb = size_bytes // (1024 * 1024)
            return False, f"Файл слишком большой ({mb} МБ). Максимум — 50 МБ."
        if self.detect_type(filename) == "unknown":
            ext = os.path.splitext(filename)[1]
            return False, (
                f"Формат {ext} не поддерживается. "
                "Принимаю: DOCX, PDF, PPTX, TXT, MP4, MP3, MOV, M4A."
            )
        return True, ""

    def extract_text(self, file_path: str) -> str:
        ext = os.path.splitext(file_path.lower())[1]
        try:
            if ext == ".txt":
                return self._read_txt(file_path)
            elif ext == ".docx":
                return self._read_docx(file_path)
            elif ext == ".pdf":
                return self._read_pdf(file_path)
            elif ext == ".pptx":
                return self._read_pptx(file_path)
            else:
                return ""
        except Exception as e:
            return f"[Ошибка чтения файла: {e}]"

    def _read_txt(self, path: str) -> str:
        with open(path, encoding="utf-8", errors="replace") as f:
            return f.read()

    def _read_docx(self, path: str) -> str:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _read_pdf(self, path: str) -> str:
        import pdfplumber
        texts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
        return "\n".join(texts)

    def _read_pptx(self, path: str) -> str:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
